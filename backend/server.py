from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel
import os
import re
import tempfile
import uuid

app = FastAPI(title="StudyBot API")

# Updated origins to match your live site
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:5173",
        "https://studybot-two.vercel.app", 
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

SESSION_STORE: dict[str, list[dict]] = {}
PDF_STORE: dict[str, str] = {}


def extract_pdf(path: str, doc_name: str) -> list[dict]:
    import fitz
    doc = fitz.open(path)
    chunks = []
    for i, page in enumerate(doc):
        text = page.get_text().strip()
        if text:
            chunks.append({
                "id": str(uuid.uuid4()),
                "doc_name": doc_name,
                "page": i + 1,
                "slide": None,
                "text": text,
            })
    return chunks


def extract_pptx(path: str, doc_name: str) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(path)
    chunks = []
    for i, slide in enumerate(prs.slides):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        combined = "\n".join(texts)
        if combined:
            chunks.append({
                "id": str(uuid.uuid4()),
                "doc_name": doc_name,
                "page": None,
                "slide": i + 1,
                "text": combined,
            })
    return chunks


STOPWORDS = {
    "a","an","the","is","are","was","were","be","been","being",
    "do","does","did","will","would","could","should","can","shall",
    "have","has","had","i","you","we","they","he","she","it",
    "in","on","at","to","for","of","and","or","but","not",
    "what","how","why","when","where","which","who","this","that",
    "with","from","by","as","into","about","also","its","if",
}

def tokenize(text: str) -> set:
    words = re.findall(r"\w+", text.lower())
    return {w for w in words if w not in STOPWORDS and len(w) > 2}


YES_NO_TRIGGERS = (
    "is ","are ","does ","do ","can ","could ","will ","would ",
    "has ","have ","did ","was ","were ","should ","shall ",
)

def classify_question(question: str) -> str:
    q = question.strip().lower()
    if any(q.startswith(t) for t in YES_NO_TRIGGERS):
        return "yes_no"
    return "open"


def retrieve(query: str, chunks: list, top_k: int = 5) -> list:
    query_tokens = tokenize(query)
    scored = []
    for chunk in chunks:
        score = len(query_tokens & tokenize(chunk["text"]))
        if score > 0:
            scored.append((score, chunk))
    scored.sort(key=lambda x: x[0], reverse=True)
    return [c for _, c in scored[:top_k]]


def extract_relevant_sentences(chunk_text: str, query_tokens: set, min_score: int = 1) -> list:
    raw = re.split(r'(?<=[.!?])\s+|\n', chunk_text.strip())
    results = []
    for sent in raw:
        sent = sent.strip()
        if len(sent) < 20:
            continue
        score = len(tokenize(sent) & query_tokens)
        if score >= min_score:
            results.append((score, sent))
    results.sort(key=lambda x: x[0], reverse=True)
    return [s for _, s in results[:3]]


POSITIVE_KEYWORDS = {
    "yes","can","does","is","are","will","helps","help","enables",
    "allows","supports","improves","used","uses","applies","works",
    "effective","useful","important","possible","advantage","benefit",
    "capable","able","achieve","solves","solve","provides","gives",
    "defined","definition","means","refers","known","called","type",
    "method","process","technique","algorithm","approach","system",
}

NEGATIVE_KEYWORDS = {
    "no","not","cannot","cant","doesnt","dont","wont","unable",
    "fails","fail","impossible","never","neither","nor","lack",
    "limitation","disadvantage","ineffective","incorrect","wrong",
    "false","unrelated","irrelevant","opposite","unlike","instead",
}

def decide_verdict(chunks: list, query_tokens: set) -> str:
    evidence_for = False
    evidence_against = False
    for chunk in chunks:
        words = set(re.findall(r"\w+", chunk["text"].lower()))
        if words & POSITIVE_KEYWORDS:
            evidence_for = True
        if words & NEGATIVE_KEYWORDS:
            evidence_against = True
    if evidence_against and not evidence_for:
        return "no"
    return "yes"


def build_correction(chunks: list, query_tokens: set) -> str:
    definition_patterns = [
        r"(?:is defined as|is called|refers to|means|is a|are a|known as)[^.!\n]{10,200}[.!\n]",
        r"(?:definition|Definition)[:\s]+[^.!\n]{10,200}[.!\n]",
    ]
    for chunk in chunks:
        for pattern in definition_patterns:
            matches = re.findall(pattern, chunk["text"], re.IGNORECASE)
            for match in matches:
                match = match.strip()
                if len(match) > 30 and len(tokenize(match) & query_tokens) > 0:
                    return match
    for chunk in chunks:
        sents = extract_relevant_sentences(chunk["text"], query_tokens, min_score=1)
        if sents:
            return sents[0]
    return ""


def build_answer(question: str, chunks: list) -> dict:
    q_type = classify_question(question)
    query_tokens = tokenize(question)

    verdict = decide_verdict(chunks, query_tokens) if q_type == "yes_no" else "answered"

    # collect proof sentences from all top chunks
    proof_sentences = []
    for chunk in chunks[:4]:
        loc = f"Slide {chunk['slide']}" if chunk["slide"] else f"Page {chunk['page']}"
        sents = extract_relevant_sentences(chunk["text"], query_tokens, min_score=1)
        for sent in sents:
            proof_sentences.append({
                "doc_name": chunk["doc_name"],
                "location": loc,
                "excerpt": sent,
                "page": chunk["page"],
                "slide": chunk["slide"],
            })

    # build explanation
    if verdict == "no":
        correct = build_correction(chunks, query_tokens)
        if correct:
            explanation = (
                f"No, that statement is incorrect. "
                f"According to the document: {correct.strip().rstrip('.')}."
            )
        else:
            explanation = "No, that statement does not match what is described in the document."

    elif verdict == "yes":
        if proof_sentences:
            parts = [p["excerpt"].strip() for p in proof_sentences[:2]]
            explanation = "Yes. " + " ".join(parts)
        else:
            explanation = "Yes, this is supported by the document."

    else:
        if proof_sentences:
            explanation = " ".join(p["excerpt"].strip() for p in proof_sentences[:3])
        else:
            explanation = "Relevant content found but no clear sentences could be extracted."

    # deduplicate sources
    seen = set()
    sources = []
    for p in proof_sentences:
        key = f"{p['doc_name']}::{p['location']}::{p['excerpt'][:40]}"
        if key not in seen:
            seen.add(key)
            sources.append({
                "doc_name": p["doc_name"],
                "location": p["location"],
                "excerpt": p["excerpt"],
                "page": p["page"],
                "slide": p["slide"],
            })

    return {
        "question_type": q_type,
        "verdict": verdict,
        "explanation": explanation,
        "sources": sources,
    }


@app.post("/session")
def create_session():
    sid = str(uuid.uuid4())
    SESSION_STORE[sid] = []
    return {"session_id": sid}


@app.post("/upload/{session_id}")
async def upload_document(session_id: str, file: UploadFile = File(...)):
    if session_id not in SESSION_STORE:
        raise HTTPException(status_code=404, detail="Session not found")

    suffix = os.path.splitext(file.filename)[1].lower()
    if suffix not in (".pdf", ".pptx", ".ppt"):
        raise HTTPException(status_code=400, detail="Only PDF and PPTX supported")

    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(await file.read())
        tmp_path = tmp.name

    if suffix == ".pdf":
        chunks = extract_pdf(tmp_path, file.filename)
        PDF_STORE[f"{session_id}::{file.filename}"] = tmp_path
    else:
        chunks = extract_pptx(tmp_path, file.filename)
        os.unlink(tmp_path)

    SESSION_STORE[session_id].extend(chunks)
    return {
        "doc_name": file.filename,
        "chunks_added": len(chunks),
        "total_chunks": len(SESSION_STORE[session_id]),
    }


@app.get("/documents/{session_id}")
def list_documents(session_id: str):
    if session_id not in SESSION_STORE:
        raise HTTPException(status_code=404, detail="Session not found")
    docs: dict[str, int] = {}
    for ch in SESSION_STORE[session_id]:
        docs[ch["doc_name"]] = docs.get(ch["doc_name"], 0) + 1
    return {"documents": [{"name": k, "chunks": v} for k, v in docs.items()]}


class AskRequest(BaseModel):
    question: str


@app.post("/ask/{session_id}")
def ask(session_id: str, body: AskRequest):
    if session_id not in SESSION_STORE:
        raise HTTPException(status_code=404, detail="Session not found")

    chunks = SESSION_STORE[session_id]
    if not chunks:
        raise HTTPException(status_code=400, detail="No documents uploaded yet")

    relevant = retrieve(body.question, chunks, top_k=5)

    if not relevant:
        return {
            "question_type": classify_question(body.question),
            "verdict": "not_found",
            "explanation": "This topic is not covered in the uploaded documents.",
            "sources": [],
        }

    return build_answer(body.question, relevant)


@app.get("/view-doc/{session_id}/{doc_name}")
def view_doc(session_id: str, doc_name: str):
    key = f"{session_id}::{doc_name}"
    path = PDF_STORE.get(key)
    if not path or not os.path.exists(path):
        raise HTTPException(status_code=404, detail="Document not available for viewing")
    return FileResponse(path, media_type="application/pdf")